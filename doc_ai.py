import io, base64
from typing import List, Dict, Optional
from PIL import Image
from openai import OpenAI
from docx import Document
from pptx import Presentation
import fitz
from pathlib import Path



class DocAI:
    def __init__(self, client: OpenAI):
        self.client = client
        self.kb: List[Dict[str, str]] = []

    def img64_from_bytes(self, img_bytes: bytes) -> Optional[str]:
        try:
            img = Image.open(io.BytesIO(img_bytes))
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            data = base64.b64encode(buffered.getvalue()).decode()
            return f"data:image/png;base64,{data}"
        except Exception as e:
            print(f"Error converting image: {e}")
            return None

    def analyze_img_from_bytes(self, img_bytes: bytes, name: str, context: str = "") -> str:
        img_url = self.img64_from_bytes(img_bytes)
        if not img_url:
            return f"Could not analyze image bytes for {name}"
        r = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": f"Describe all details from image in {name}. {context}".strip()},
                    {"type": "image_url", "image_url": {"url": img_url}}
                ]
            }],
            max_tokens=1000
        )
        return r.choices[0].message.content
    
    def text_docx_from_bytes(self, name: str, file_bytes: bytes) -> str:
        doc = Document(io.BytesIO(file_bytes))
        content_parts: List[str] = []

        # Text
        text_content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        if text_content:
            content_parts.append(f"TEXT CONTENT:\n{text_content}")

        # Images
        img_count = 0
        for rel in doc.part.rels.values():
            if "image" in getattr(rel, "target_ref", ""):
                try:
                    img_bytes = rel.target_part.blob
                    img_desc = self.analyze_img_from_bytes(
                        img_bytes,
                        name,
                        f"This is image {img_count + 1} from the document"
                    )
                    img_count += 1
                    content_parts.append(f"\nIMAGE {img_count} DESCRIPTION:\n{img_desc}")
                except Exception as e:
                    print(f"Error processing image: {e}")

        return "\n\n".join(content_parts) if content_parts else "No content found"

    def text_pptx_from_bytes(self, name: str, file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        content_parts: List[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # Text
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                content_parts.append(f"SLIDE {slide_num} TEXT:\n" + "\n".join(slide_text))

            # Images
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        img_desc = self.analyze_img_from_bytes(
                            img_bytes,
                            name,
                            f"Image from slide {slide_num}"
                        )
                        content_parts.append(f"SLIDE {slide_num} IMAGE:\n{img_desc}")
                    except Exception as e:
                        print(f"Error processing image from slide {slide_num}: {e}")

        return "\n\n".join(content_parts) if content_parts else "No content found"

    def text_pdf_from_bytes(self, name: str, file_bytes: bytes) -> str:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        out: List[str] = []
        for i, pg in enumerate(doc):
            # Render page as image
            pix = pg.get_pixmap(matrix=fitz.Matrix(2, 2))
            png_bytes = pix.tobytes("png")
            b64 = base64.b64encode(png_bytes).decode()

            r = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": f"Extract all content (text and images) from page {i+1} of {name}"},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}}
                    ]
                }]
            )
            out.append(f"PAGE {i+1}:\n{r.choices[0].message.content}")
        return "\n\n".join(out)
    
    def add_attachment(self, filename: str, content_type: str, file_bytes: bytes) -> str:
        """
        Process and add an attachment to KB. Returns a status message.
        """
        ext = Path(filename).suffix.lower()
        print(f"Processing {filename}...")

        if ext in [".png", ".jpg", ".jpeg"]:
            desc = self.analyze_img_from_bytes(file_bytes, name=filename)
            content = desc
        elif ext == ".pdf":
            content = self.text_pdf_from_bytes(filename, file_bytes)
        elif ext == ".docx":
            content = self.text_docx_from_bytes(filename, file_bytes)
        elif ext in [".pptx", ".ppt"]:
            content = self.text_pptx_from_bytes(filename, file_bytes)
        else:
            return f"Unsupported type: {ext or content_type}"

        return content

    def build_context(self) -> Optional[str]:
        if not self.kb:
            return None
        return "\n\n".join([f"=== {d['name']} ===\n{d['content']}" for d in self.kb])


